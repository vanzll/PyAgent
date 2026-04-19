from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).with_name("cs5260_financial_agent_demo.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.author = "Zhenglin Wan"
prs.core_properties.title = "PyCallingAgent Research Desk"
prs.core_properties.subject = "Financial research agent for CS5260"

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"

COLORS = {
    "navy": RGBColor(0x0B, 0x14, 0x23),
    "navy_soft": RGBColor(0x14, 0x24, 0x3C),
    "navy_panel": RGBColor(0x10, 0x21, 0x39),
    "cream": RGBColor(0xF7, 0xF4, 0xED),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "sand": RGBColor(0xEE, 0xE5, 0xD6),
    "ink": RGBColor(0x16, 0x20, 0x33),
    "ink_soft": RGBColor(0x5E, 0x6B, 0x7E),
    "gold": RGBColor(0xD2, 0xA4, 0x5C),
    "gold_soft": RGBColor(0xF0, 0xCF, 0x93),
    "blue": RGBColor(0x4E, 0x7D, 0xD1),
    "green": RGBColor(0x2B, 0x8C, 0x74),
    "red": RGBColor(0xB9, 0x5C, 0x5C),
    "border": RGBColor(0xD8, 0xDD, 0xE6),
    "blue_bg": RGBColor(0xEA, 0xF0, 0xF9),
    "green_bg": RGBColor(0xEE, 0xF8, 0xF5),
    "red_bg": RGBColor(0xFC, 0xEF, 0xEF),
    "gold_bg": RGBColor(0xFB, 0xF3, 0xE5),
}

slide_no = 0


def add_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_card(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_text(slide, x, y, w, h, text, *, font_name=BODY_FONT, size=14, color=None, bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_pill(slide, x, y, w, text, fill, color):
    add_card(slide, x, y, w, 0.34, fill, fill)
    add_text(slide, x, y + 0.04, w, 0.2, text, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_bullets(slide, x, y, w, items, *, size=13, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(2.6))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(8)
    return box


def add_slide_number(slide, dark=False):
    global slide_no
    slide_no += 1
    add_text(
        slide,
        12.48,
        7.0,
        0.3,
        0.18,
        str(slide_no),
        size=10,
        color=RGBColor(0xC9, 0xD3, 0xE3) if dark else RGBColor(0x7C, 0x87, 0x97),
        align=PP_ALIGN.RIGHT,
    )


def add_title(slide, title: str, subtitle: str) -> None:
    add_text(slide, 0.55, 0.42, 8.8, 0.42, title, font_name=TITLE_FONT, size=26, color=COLORS["ink"], bold=True)
    add_text(slide, 0.55, 1.0, 8.7, 0.24, subtitle, size=11, color=COLORS["ink_soft"])


def add_dark_title(slide, title: str, subtitle: str) -> None:
    add_text(slide, 0.55, 0.48, 8.8, 0.82, title, font_name=TITLE_FONT, size=26, color=COLORS["white"], bold=True)
    add_text(slide, 0.55, 1.3, 8.5, 0.24, subtitle, size=12, color=RGBColor(0xD7, 0xDF, 0xEC))


def add_footer(slide, text: str, *, dark=False) -> None:
    add_text(
        slide,
        0.55,
        6.92,
        11.6,
        0.16,
        text,
        size=9,
        color=RGBColor(0xB9, 0xC6, 0xD8) if dark else RGBColor(0x7A, 0x87, 0x97),
        italic=True,
    )


def build_title_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])
    add_pill(slide, 0.55, 0.62, 2.4, "NUS CS5260 FINAL PROJECT", COLORS["gold"], COLORS["ink"])
    add_text(slide, 0.55, 1.45, 8.2, 0.6, "PyCallingAgent Research Desk", font_name=TITLE_FONT, size=31, color=COLORS["white"], bold=True)
    add_text(slide, 0.55, 2.2, 7.8, 0.26, "A public financial research agent for U.S. stocks and ETFs", size=15, color=RGBColor(0xD7, 0xDF, 0xEC))
    add_text(
        slide,
        0.55,
        3.0,
        6.1,
        1.2,
        "The product accepts tickers and a research question, then fetches public market data, filing-backed facts, comparison tables, and chart artifacts through a PyCallingAgent-powered runtime.",
        size=17,
        color=COLORS["cream"],
    )

    metric_positions = [
        (7.8, 1.56, 1.55, "3", "Public data sources"),
        (9.52, 1.56, 1.55, "SSE", "Live run stream"),
        (11.24, 1.56, 1.55, "0 SGD", "Demo-mode cost"),
        (7.8, 2.78, 2.35, "3 outputs", "exported"),
        (10.32, 2.78, 2.47, "Adaptive", "task-driven output"),
    ]
    for x, y, w, value, label in metric_positions:
        add_card(slide, x, y, w, 1.0, COLORS["navy_panel"])
        add_text(slide, x + 0.14, y + 0.12, w - 0.2, 0.22, value, font_name=TITLE_FONT, size=20, color=COLORS["gold_soft"], bold=True)
        add_text(slide, x + 0.14, y + 0.56, w - 0.2, 0.14, label, size=9, color=RGBColor(0xD3, 0xDC, 0xEC))

    add_card(slide, 7.8, 4.15, 5.0, 1.6, COLORS["navy_panel"])
    add_text(slide, 8.02, 4.36, 2.1, 0.18, "Public homepage statement", size=11, color=COLORS["gold_soft"], bold=True)
    add_text(
        slide,
        8.02,
        4.66,
        4.5,
        0.6,
        "This public-facing agent product is entirely derived from work conducted as part of the NUS CS5260 course.",
        size=13,
        color=COLORS["white"],
    )
    add_text(slide, 0.55, 6.52, 3.2, 0.16, "Zhenglin Wan · April 17, 2026", size=11, color=RGBColor(0xC9, 0xD3, 0xE3))
    add_slide_number(slide, dark=True)


def build_problem_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["cream"])
    add_title(slide, "Why this product exists", "Financial research is easy to start and tedious to finish.")

    add_card(slide, 0.62, 1.58, 5.96, 4.95, COLORS["white"])
    add_pill(slide, 0.9, 1.84, 1.46, "ANALYST PAIN", COLORS["blue_bg"], COLORS["blue"])
    add_text(slide, 0.9, 2.22, 4.6, 0.54, "Manual workflows break the research loop", font_name=TITLE_FONT, size=19, color=COLORS["ink"], bold=True)
    add_bullets(
        slide,
        0.95,
        3.08,
        5.0,
        [
            "Price data, filings, and macro context live in different tools.",
            "A simple comparison request still requires multiple tabs and export steps.",
            "Outputs become inconsistent: screenshots, notes, and ad hoc charts.",
            "For a live course demo, external API or model failures can kill momentum.",
        ],
        size=15,
    )

    add_card(slide, 6.78, 1.58, 5.93, 4.95, COLORS["navy_soft"])
    add_pill(slide, 7.05, 1.84, 1.7, "DESIGN RESPONSE", RGBColor(0x24, 0x3A, 0x60), COLORS["gold_soft"])
    add_text(slide, 7.05, 2.22, 4.9, 0.54, "A public agent workbench solves the end-to-end task", font_name=TITLE_FONT, size=19, color=COLORS["white"], bold=True)
    add_bullets(
        slide,
        7.1,
        3.08,
        5.0,
        [
            "Input is natural language plus one or more tickers.",
            "The agent chooses a skill path and grounds outputs in public data.",
            "The UI streams progress and exports summary, table, and chart artifacts.",
            "Demo mode uses deterministic examples, so the classroom demo remains stable.",
        ],
        size=15,
        color=RGBColor(0xF3, 0xF6, 0xFA),
    )
    add_footer(slide, "The scope is intentionally narrow: research assistant, not investment advisor.")
    add_slide_number(slide)


def build_overview_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["cream"])
    add_title(slide, "Product overview", "One question in, three artifact types out.")
    add_card(slide, 0.72, 1.55, 12.0, 5.25, RGBColor(0xFB, 0xFA, 0xF7))

    add_text(slide, 1.02, 1.95, 1.0, 0.16, "Input", size=12, color=COLORS["ink_soft"], bold=True)
    add_card(slide, 1.0, 2.24, 3.18, 3.78, COLORS["white"])
    add_text(slide, 1.24, 2.5, 2.4, 0.48, "Tickers and question", font_name=TITLE_FONT, size=17, color=COLORS["ink"], bold=True)
    add_text(slide, 1.24, 3.04, 1.4, 0.16, "Examples", size=11, color=COLORS["blue"], bold=True)
    add_bullets(
        slide,
        1.24,
        3.32,
        2.52,
        [
            "AAPL: summarize the company and show recent trend",
            "AMD, NVDA: compare valuation, growth, and performance",
            "SPY: give me a market snapshot and the main risks",
        ],
        size=13,
    )

    add_text(slide, 4.92, 1.95, 1.4, 0.16, "Agent runtime", size=12, color=COLORS["ink_soft"], bold=True)
    add_card(slide, 4.9, 2.24, 3.52, 3.78, COLORS["navy_soft"])
    steps = [
        "01  Fetch market data",
        "02  Normalize into one workspace",
        "03  Call comparison and chart skills",
        "04  Return research brief + artifacts",
    ]
    for idx, step in enumerate(steps):
        y = 2.52 + idx * 0.78
        add_card(slide, 5.18, y, 2.94, 0.56, RGBColor(0x1B, 0x31, 0x50))
        add_text(slide, 5.32, y + 0.12, 2.54, 0.16, step, size=11.5, color=COLORS["white"])

    add_text(slide, 9.12, 1.95, 1.2, 0.16, "Outputs", size=12, color=COLORS["ink_soft"], bold=True)
    outputs = [
        (2.24, COLORS["gold_bg"], RGBColor(0xEF, 0xD7, 0xAB), "Markdown brief"),
        (3.48, RGBColor(0xEE, 0xF4, 0xFB), RGBColor(0xD7, 0xE2, 0xF2), "Comparison table"),
        (4.72, COLORS["green_bg"], RGBColor(0xD7, 0xEC, 0xE5), "Chart artifact"),
    ]
    for y, fill, line, label in outputs:
        add_card(slide, 9.08, y, 2.94, 1.08, fill, line)
        add_text(slide, 9.34, y + 0.38, 1.8, 0.18, label, font_name=TITLE_FONT, size=16, color=COLORS["ink"], bold=True)

    add_text(slide, 1.04, 6.2, 10.8, 0.16, "The app behaves like a small research desk rather than a static market dashboard.", size=11, color=COLORS["ink_soft"], italic=True)
    add_slide_number(slide)


def build_architecture_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["cream"])
    add_title(slide, "System architecture", "Structured data providers on the outside, PyCallingAgent runtime in the middle.")

    add_card(slide, 0.78, 1.78, 2.25, 4.7, COLORS["white"])
    add_pill(slide, 1.0, 2.02, 1.56, "DATA PROVIDERS", COLORS["blue_bg"], COLORS["blue"])
    providers = [
        ("Alpha Vantage", "Prices, overview, optional news"),
        ("SEC EDGAR", "Recent filings and company facts"),
        ("FRED", "Macro rates and inflation series"),
    ]
    for idx, (name, desc) in enumerate(providers):
        y = 2.56 + idx * 1.12
        add_card(slide, 1.02, y, 1.78, 0.94, RGBColor(0xFB, 0xFA, 0xF7))
        add_text(slide, 1.16, y + 0.14, 1.4, 0.24, name, font_name=TITLE_FONT, size=13, color=COLORS["ink"], bold=True)
        add_text(slide, 1.16, y + 0.48, 1.42, 0.18, desc, size=9.5, color=COLORS["ink_soft"])

    add_text(slide, 3.24, 3.74, 0.4, 0.24, "→", size=26, color=COLORS["gold"], bold=True)

    add_card(slide, 3.72, 1.78, 5.24, 4.7, COLORS["navy_soft"])
    add_pill(slide, 4.02, 2.02, 1.9, "PYCALLINGAGENT RUNTIME", RGBColor(0x24, 0x3A, 0x60), COLORS["gold_soft"])
    add_text(slide, 4.02, 2.46, 3.8, 0.24, "Unified research workspace", font_name=TITLE_FONT, size=18, color=COLORS["white"], bold=True)
    add_text(
        slide,
        4.02,
        3.06,
        4.4,
        0.5,
        "Provider outputs become typed runtime objects: price history, comparison frames, macro series, and company facts. The agent reasons over structured inputs before choosing a skill path.",
        size=11.5,
        color=RGBColor(0xE4, 0xEB, 0xF6),
    )
    skills = [
        ("market-data", "Price inspection"),
        ("fundamentals", "Filing-backed facts"),
        ("comparison", "Peer tables"),
        ("charting", "PNG outputs"),
    ]
    for idx, (name, desc) in enumerate(skills):
        x = 4.02 + (idx % 2) * 2.2
        y = 4.06 + (idx // 2) * 0.96
        add_card(slide, x, y, 2.05, 0.72, RGBColor(0x1B, 0x31, 0x50))
        add_text(slide, x + 0.14, y + 0.16, 1.76, 0.16, name, size=11, color=COLORS["gold_soft"], bold=True)
        add_text(slide, x + 0.14, y + 0.38, 1.76, 0.16, desc, size=10, color=RGBColor(0xD8, 0xE1, 0xEE))

    add_text(slide, 9.1, 3.74, 0.4, 0.24, "→", size=26, color=COLORS["gold"], bold=True)

    add_card(slide, 9.56, 1.78, 2.96, 4.7, COLORS["white"])
    add_pill(slide, 9.8, 2.02, 1.12, "WEB OUTPUT", COLORS["green_bg"], COLORS["green"])
    outputs = ["Live event stream", "Snapshot cards", "Research brief", "Downloadable artifacts"]
    for idx, label in enumerate(outputs):
        y = 2.62 + idx * 0.88
        add_card(slide, 9.82, y, 2.42, 0.58, RGBColor(0xFB, 0xFA, 0xF7))
        add_text(slide, 10.0, y + 0.16, 2.0, 0.16, label, size=11, color=COLORS["ink"])

    add_footer(slide, "The key engineering choice is normalizing provider output before the agent starts reasoning.")
    add_slide_number(slide)


def build_trust_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["cream"])
    add_title(slide, "Grounding, safety, and cost control", "The project stays useful by staying narrow.")

    cards = [
        (0.74, "Grounded outputs", ("PUBLIC DATA", COLORS["blue_bg"], COLORS["blue"]), [
            "Price history and profile data come from Alpha Vantage.",
            "SEC EDGAR provides official filings and company facts.",
            "FRED adds macro context when the question needs it.",
        ]),
        (4.44, "Safe product scope", ("RESEARCH ONLY", COLORS["red_bg"], COLORS["red"]), [
            "The app is framed as a financial research assistant.",
            "Every run is labeled: For research use only. Not investment advice.",
            "No trading execution, portfolio optimization, or personal finance claims.",
        ]),
        (8.14, "Low-risk demo path", ("DEMO MODE", COLORS["green_bg"], COLORS["green"]), [
            "Deterministic market examples remove rate-limit risk.",
            "Demo mode can run with zero API and model cost.",
            "The same UI, outputs, and artifact flow are preserved in class.",
        ]),
    ]
    for x, title, pill, lines in cards:
        add_card(slide, x, 1.82, 3.2, 4.52, COLORS["white"])
        add_pill(slide, x + 0.26, 2.08, 1.3, pill[0], pill[1], pill[2])
        add_text(slide, x + 0.26, 2.46, 2.5, 0.48, title, font_name=TITLE_FONT, size=18, color=COLORS["ink"], bold=True)
        add_bullets(slide, x + 0.28, 3.18, 2.5, lines, size=13)

    add_card(slide, 0.92, 6.4, 11.54, 0.48, COLORS["gold_bg"], RGBColor(0xF1, 0xD7, 0xA7))
    add_text(
        slide,
        1.16,
        6.54,
        10.9,
        0.16,
        "Homepage statement: This public-facing agent product is entirely derived from work conducted as part of the NUS CS5260 course.",
        size=11,
        color=COLORS["ink"],
    )
    add_footer(slide, "Official sources: Alpha Vantage documentation, SEC EDGAR API, and FRED API.")
    add_slide_number(slide)


def build_demo_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["cream"])
    add_title(slide, "5-minute demo plan", "Three runs, each chosen to show a different product capability.")

    demos = [
        (0.78, "RUN 01", COLORS["blue_bg"], COLORS["blue"], "AAPL single-stock brief", "Summarize AAPL and show the recent price trend.", ["Snapshot card", "Trend chart", "Markdown brief"]),
        (4.46, "RUN 02", COLORS["gold_bg"], COLORS["gold"], "AMD vs NVDA comparison", "Compare valuation, latest fundamentals, and recent performance.", ["Two snapshot cards", "Comparison CSV", "Relative returns chart"]),
        (8.14, "RUN 03", COLORS["green_bg"], COLORS["green"], "SPY market snapshot", "Give me a market snapshot and highlight the main risks.", ["ETF framing", "Macro context", "Short risk summary"]),
    ]
    for x, tag, bg, color, title, prompt, points in demos:
        add_card(slide, x, 1.86, 3.18, 4.66, COLORS["white"])
        add_pill(slide, x + 0.26, 2.1, 0.92, tag, bg, color)
        add_text(slide, x + 0.26, 2.48, 2.42, 0.36, title, font_name=TITLE_FONT, size=18, color=COLORS["ink"], bold=True)
        add_text(slide, x + 0.26, 3.02, 2.46, 0.66, prompt, size=12, color=COLORS["ink_soft"], italic=True)
        add_text(slide, x + 0.26, 3.9, 1.1, 0.16, "Show", size=11, color=color, bold=True)
        add_bullets(slide, x + 0.28, 4.18, 2.3, points, size=13)

    add_footer(slide, "Use deterministic demo mode in class to avoid API or model latency surprises.")
    add_slide_number(slide)


def build_evaluation_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["cream"])
    add_title(slide, "Why this is an agent, not a dashboard", "The output changes based on the task, not just the selected ticker.")

    add_card(slide, 0.9, 1.82, 11.55, 4.8, COLORS["white"])
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(2.16), Inches(11.55), Inches(0.52))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["navy_soft"]
    header.line.color.rgb = COLORS["navy_soft"]
    add_text(slide, 1.16, 2.32, 2.0, 0.16, "Criterion", size=11, color=COLORS["white"], bold=True)
    add_text(slide, 4.18, 2.32, 5.0, 0.16, "How the project satisfies it", size=11, color=COLORS["white"], bold=True)

    rows = [
        ("Open-ended task interpretation", "The user provides a natural-language research prompt, not a fixed dashboard control."),
        ("Skill selection", "The runtime can route through market-data, fundamentals, comparison, and charting skills depending on the question."),
        ("Artifact generation", "The agent returns a brief, a structured table artifact, and a chart artifact in one run."),
        ("Grounded evidence", "SEC EDGAR and structured provider objects are injected before reasoning starts."),
        ("Presentation reliability", "Demo mode preserves the exact agent workflow while removing rate-limit and cost risk."),
    ]
    for idx, (criterion, detail) in enumerate(rows):
        y = 2.68 + idx * 0.76
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.02), Inches(y), Inches(11.2), Inches(0.01))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE7, 0xEB, 0xF1)
        line.line.color.rgb = RGBColor(0xE7, 0xEB, 0xF1)
        add_text(slide, 1.16, y + 0.18, 2.6, 0.18, criterion, font_name=TITLE_FONT, size=14, color=COLORS["ink"], bold=True)
        add_text(slide, 4.18, y + 0.14, 7.55, 0.32, detail, size=12, color=COLORS["ink_soft"])

    add_footer(slide, "A static dashboard cannot adapt its output structure to three different prompts this way.")
    add_slide_number(slide)


def build_closing_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])
    add_dark_title(slide, "What is complete, and\nwhat comes next", "The MVP is ready for a stable classroom demo and public deployment.")

    sections = [
        (0.76, "READY NOW", [
            "Finance-themed web UI with course statement on the homepage",
            "Public-data provider layer with deterministic demo fallback",
            "Snapshot cards, research brief, tables, charts, and downloadable artifacts",
            "Three rehearsed demo scenarios for the final presentation",
        ]),
        (5.02, "NEXT STEP", [
            "Deploy the FastAPI app to a public URL before April 17, 2026",
            "Record a 60-90 second demo video as backup",
            "Swap in live API keys for the final hosted version",
            "Optionally add news explanation after the core demo is stable",
        ]),
    ]
    for x, tag, items in sections:
        add_card(slide, x, 1.86, 4.02, 4.7, COLORS["navy_panel"])
        add_pill(slide, x + 0.26, 2.12, 1.06 if tag == "NEXT STEP" else 1.1, tag, RGBColor(0x24, 0x3A, 0x60), COLORS["gold_soft"])
        add_bullets(slide, x + 0.26, 2.56, 3.3, items, size=13, color=RGBColor(0xF4, 0xF7, 0xFB))

    add_card(slide, 9.28, 1.86, 3.22, 4.7, RGBColor(0x13, 0x27, 0x40))
    add_text(slide, 9.54, 2.18, 2.1, 0.24, "Takeaway", font_name=TITLE_FONT, size=20, color=COLORS["white"], bold=True)
    add_text(slide, 9.54, 2.72, 2.48, 1.04, "PyCallingAgent Research Desk turns a vague finance question into a grounded research packet in one run.", size=15, color=RGBColor(0xE2, 0xEA, 0xF7))
    add_text(slide, 9.54, 4.18, 2.5, 0.26, "For research use only. Not investment advice.", size=13, color=COLORS["gold_soft"], bold=True)
    add_text(slide, 9.54, 5.1, 2.5, 0.5, "Public URL placeholder:\nhttps://your-deployment-url", size=11, color=RGBColor(0xC9, 0xD3, 0xE3))
    add_footer(slide, "This public-facing agent product is entirely derived from work conducted as part of the NUS CS5260 course.", dark=True)
    add_slide_number(slide, dark=True)


def main() -> None:
    build_title_slide()
    build_problem_slide()
    build_overview_slide()
    build_architecture_slide()
    build_trust_slide()
    build_demo_slide()
    build_evaluation_slide()
    build_closing_slide()
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
