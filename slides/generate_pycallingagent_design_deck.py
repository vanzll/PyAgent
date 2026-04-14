from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).with_name("pycallingagent_design_deployment.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.author = "Zhenglin Wan"
prs.core_properties.title = "PyCallingAgent Design and Deployment"
prs.core_properties.subject = "CS5260 final project presentation"

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"

COLORS = {
    "navy": RGBColor(0x0D, 0x1B, 0x2A),
    "navy_soft": RGBColor(0x1B, 0x26, 0x3A),
    "teal": RGBColor(0x0F, 0x5C, 0x52),
    "teal_soft": RGBColor(0xD9, 0xEE, 0xEA),
    "sand": RGBColor(0xF4, 0xEF, 0xE5),
    "paper": RGBColor(0xFB, 0xF8, 0xF1),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "ink": RGBColor(0x1F, 0x23, 0x28),
    "muted": RGBColor(0x67, 0x5F, 0x56),
    "gold": RGBColor(0xD2, 0xA4, 0x5C),
    "gold_soft": RGBColor(0xF1, 0xD3, 0x9D),
    "line": RGBColor(0xD7, 0xCF, 0xC2),
    "blue_bg": RGBColor(0xEA, 0xF0, 0xF9),
    "green_bg": RGBColor(0xEC, 0xF7, 0xF4),
    "amber_bg": RGBColor(0xFB, 0xF3, 0xE4),
    "rose_bg": RGBColor(0xFB, 0xEC, 0xEA),
}

slide_no = 0


def add_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font_name=BODY_FONT,
    size=14,
    color=None,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_card(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_pill(slide, x, y, w, text, fill, color):
    add_card(slide, x, y, w, 0.34, fill, fill)
    add_text(slide, x, y + 0.04, w, 0.18, text, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_title(slide, title: str, subtitle: str, *, dark=False) -> None:
    color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = COLORS["gold_soft"] if dark else COLORS["muted"]
    add_text(slide, 0.6, 0.45, 8.8, 0.44, title, font_name=TITLE_FONT, size=27, color=color, bold=True)
    add_text(slide, 0.6, 0.98, 8.8, 0.22, subtitle, size=11, color=sub_color)


def add_bullets(slide, x, y, w, items, *, size=14, color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4.2))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {item}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(8)
    return box


def add_slide_number(slide, *, dark=False) -> None:
    global slide_no
    slide_no += 1
    add_text(
        slide,
        12.55,
        7.0,
        0.25,
        0.16,
        str(slide_no),
        size=10,
        color=COLORS["gold_soft"] if dark else COLORS["muted"],
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide, text: str, *, dark=False) -> None:
    add_text(slide, 0.6, 6.95, 10.8, 0.15, text, size=9, color=COLORS["gold_soft"] if dark else COLORS["muted"], italic=True)


def add_connector(slide, x1, y1, x2, y2, color):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    return line


def build_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])
    add_pill(slide, 0.62, 0.62, 2.05, "NUS CS5260 FINAL PROJECT", COLORS["gold"], COLORS["ink"])
    add_text(slide, 0.62, 1.45, 7.9, 0.58, "PyCallingAgent", font_name=TITLE_FONT, size=31, color=COLORS["white"], bold=True)
    add_text(slide, 0.62, 2.16, 8.4, 0.24, "Prompt-first financial research agent with chat, charts, tables, and downloadable reports", size=15, color=RGBColor(0xD4, 0xDE, 0xEC))
    add_text(
        slide,
        0.62,
        3.02,
        5.9,
        1.0,
        "Goal: make the product feel like a real agent, not a static dashboard. One prompt in, grounded research outputs out.",
        size=18,
        color=COLORS["paper"],
    )

    stats = [
        ("Prompt-first", "single input"),
        ("Chat + Artifacts", "LLM answer plus outputs"),
        ("Public URL", "Render deployment"),
        ("Research only", "not investment advice"),
    ]
    x = 7.55
    for idx, (value, label) in enumerate(stats):
        col = idx % 2
        row = idx // 2
        left = x + col * 2.45
        top = 1.65 + row * 1.45
        add_card(slide, left, top, 2.1, 1.02, COLORS["navy_soft"])
        add_text(slide, left + 0.14, top + 0.12, 1.8, 0.24, value, font_name=TITLE_FONT, size=17, color=COLORS["gold_soft"], bold=True)
        add_text(slide, left + 0.14, top + 0.58, 1.8, 0.16, label, size=10, color=RGBColor(0xD4, 0xDE, 0xEC))

    add_card(slide, 7.55, 4.72, 4.95, 1.18, COLORS["navy_soft"])
    add_text(slide, 7.8, 4.95, 4.4, 0.22, "Public statement on homepage", size=11, color=COLORS["gold_soft"], bold=True)
    add_text(
        slide,
        7.8,
        5.28,
        4.25,
        0.4,
        "This public-facing agent product is entirely derived from work conducted as part of the NUS CS5260 course.",
        size=12,
        color=COLORS["white"],
    )
    add_footer(slide, "Repository: github.com/vanzll/PyAgent · Public URL: pycallingagent.onrender.com", dark=True)
    add_slide_number(slide, dark=True)


def build_design_principles_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["sand"])
    add_title(slide, "Design principles", "Why the product is shaped like an agent instead of a financial dashboard.")

    cards = [
        (0.72, 1.7, COLORS["blue_bg"], "1. Prompt-first", "The user starts with one natural-language prompt. The system infers company, ETF, or market proxy instead of forcing structured filters."),
        (3.98, 1.7, COLORS["green_bg"], "2. Artifact-oriented", "A useful agent does not stop at text. It should also produce tables, charts, and downloadable outputs in the same turn."),
        (7.24, 1.7, COLORS["amber_bg"], "3. Public-data grounded", "Outputs are grounded in public market data and SEC filings so the demo looks practical instead of purely generative."),
        (10.5, 1.7, COLORS["rose_bg"], "4. Resilient by design", "The classroom demo cannot depend on one fragile API. The system keeps working through public data mode and stable fallbacks."),
    ]

    for x, y, fill, title, body in cards:
        add_card(slide, x, y, 2.5, 3.88, fill, COLORS["line"])
        add_text(slide, x + 0.18, y + 0.22, 2.05, 0.46, title, font_name=TITLE_FONT, size=18, color=COLORS["ink"], bold=True)
        add_text(slide, x + 0.18, y + 0.94, 2.05, 2.2, body, size=13, color=COLORS["ink"])

    add_card(slide, 0.72, 5.92, 12.0, 0.62, COLORS["white"], COLORS["line"])
    add_text(
        slide,
        0.96,
        6.11,
        11.4,
        0.16,
        "Result: the UX behaves like a small analyst assistant. The user asks a question, the agent reasons over data, and the UI shows both the answer and the work product.",
        size=12,
        color=COLORS["muted"],
    )
    add_slide_number(slide)


def build_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["paper"])
    add_title(slide, "Agent architecture", "A thin web product on top of CaveAgent runtime and skills.")

    add_card(slide, 0.8, 2.0, 2.0, 1.24, COLORS["white"], COLORS["line"])
    add_text(slide, 1.05, 2.28, 1.5, 0.22, "Browser UI", font_name=TITLE_FONT, size=18, bold=True)
    add_text(slide, 1.05, 2.62, 1.55, 0.34, "Prompt input\nAssistant answer\nLive status", size=12, color=COLORS["muted"])

    add_card(slide, 3.3, 1.65, 2.38, 1.94, COLORS["blue_bg"], COLORS["line"])
    add_text(slide, 3.58, 1.98, 1.8, 0.22, "FastAPI service", font_name=TITLE_FONT, size=18, bold=True)
    add_text(slide, 3.58, 2.32, 1.9, 0.62, "Sessions\nRuns\nSSE / polling\nArtifact download", size=12, color=COLORS["muted"])

    add_card(slide, 6.2, 1.36, 3.0, 2.52, COLORS["green_bg"], COLORS["line"])
    add_text(slide, 6.5, 1.7, 2.25, 0.22, "FinancialResearchRunner", font_name=TITLE_FONT, size=18, bold=True)
    add_text(slide, 6.5, 2.02, 2.3, 0.84, "Infer tickers\nFetch public data\nInject runtime variables\nCall LLM + skills\nSave outputs", size=12, color=COLORS["muted"])

    add_card(slide, 9.7, 1.65, 2.82, 1.94, COLORS["amber_bg"], COLORS["line"])
    add_text(slide, 10.0, 1.98, 2.2, 0.22, "Data layer", font_name=TITLE_FONT, size=18, bold=True)
    add_text(slide, 10.0, 2.32, 2.2, 0.62, "SEC EDGAR\nOptional Alpha Vantage\nPublic data fallback", size=12, color=COLORS["muted"])

    add_card(slide, 3.7, 4.55, 5.4, 1.4, COLORS["white"], COLORS["line"])
    add_text(slide, 4.0, 4.88, 4.8, 0.22, "Outputs returned to the same page", font_name=TITLE_FONT, size=18, bold=True)
    add_text(slide, 4.0, 5.24, 4.8, 0.34, "Summary · comparison table · chart image · downloadable report", size=12, color=COLORS["muted"])

    add_connector(slide, 2.82, 2.62, 3.3, 2.62, COLORS["teal"])
    add_connector(slide, 5.68, 2.62, 6.2, 2.62, COLORS["teal"])
    add_connector(slide, 9.2, 2.62, 9.7, 2.62, COLORS["teal"])
    add_connector(slide, 7.7, 3.88, 7.2, 4.55, COLORS["teal"])
    add_connector(slide, 5.35, 4.55, 2.2, 3.25, COLORS["teal"])

    add_footer(slide, "The product layer stays simple. Most complexity is isolated in the runner and provider boundaries.")
    add_slide_number(slide)


def build_execution_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["sand"])
    add_title(slide, "Execution flow", "What happens after the user presses Send.")

    steps = [
        ("01", "Infer target", "Extract AAPL, NVDA, SPY, or market proxy from the prompt."),
        ("02", "Load data", "Pull public market context and SEC facts into a single normalized bundle."),
        ("03", "Run agent", "The LLM uses skills over runtime variables instead of hallucinating raw finance facts."),
        ("04", "Materialize outputs", "The turn saves charts, tables, markdown summary, and snapshot cards."),
        ("05", "Stream progress", "The page shows running status and final artifacts in one session."),
    ]

    y = 1.72
    for idx, (step, title, body) in enumerate(steps):
        fill = COLORS["white"] if idx % 2 == 0 else COLORS["green_bg"]
        add_card(slide, 1.0 + idx * 2.45, y, 2.1, 3.92, fill, COLORS["line"])
        add_text(slide, 1.22 + idx * 2.45, y + 0.22, 0.55, 0.26, step, font_name=TITLE_FONT, size=18, color=COLORS["teal"], bold=True)
        add_text(slide, 1.22 + idx * 2.45, y + 0.72, 1.55, 0.26, title, font_name=TITLE_FONT, size=17, bold=True)
        add_text(slide, 1.22 + idx * 2.45, y + 1.18, 1.6, 1.8, body, size=12, color=COLORS["muted"])
        if idx < len(steps) - 1:
            add_connector(slide, 3.08 + idx * 2.45, y + 1.96, 3.38 + idx * 2.45, y + 1.96, COLORS["teal"])

    add_card(slide, 1.0, 5.95, 11.95, 0.62, COLORS["white"], COLORS["line"])
    add_text(
        slide,
        1.24,
        6.14,
        11.3,
        0.16,
        "This is why the UI feels more agent-like than dashboard-like: each turn produces both language output and visible work products.",
        size=12,
        color=COLORS["muted"],
    )
    add_slide_number(slide)


def build_deployment_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["paper"])
    add_title(slide, "Online deployment", "The deployment path is intentionally minimal.")

    add_card(slide, 0.88, 1.75, 3.0, 4.5, COLORS["white"], COLORS["line"])
    add_text(slide, 1.16, 2.02, 2.3, 0.24, "Step 1 · GitHub", font_name=TITLE_FONT, size=19, bold=True)
    add_bullets(
        slide,
        1.16,
        2.46,
        2.3,
        [
            "Push to `main`",
            "Render watches the repo",
            "New commit triggers redeploy",
        ],
        size=13,
    )

    add_card(slide, 4.18, 1.75, 4.06, 4.5, COLORS["blue_bg"], COLORS["line"])
    add_text(slide, 4.48, 2.02, 3.2, 0.24, "Step 2 · Render Web Service", font_name=TITLE_FONT, size=19, bold=True)
    add_bullets(
        slide,
        4.48,
        2.46,
        3.2,
        [
            "Build: `pip install \".[openai,webapp]\"`",
            "Start: `python -m uvicorn cave_agent.webapp.app:create_app --factory --host 0.0.0.0 --port $PORT`",
            "Public URL comes from Render, no custom domain required",
        ],
        size=13,
    )

    add_card(slide, 8.54, 1.75, 3.88, 4.5, COLORS["green_bg"], COLORS["line"])
    add_text(slide, 8.84, 2.02, 3.0, 0.24, "Step 3 · Runtime config", font_name=TITLE_FONT, size=19, bold=True)
    add_bullets(
        slide,
        8.84,
        2.46,
        3.0,
        [
            "`LLM_MODEL_ID`",
            "`LLM_API_KEY`",
            "`LLM_BASE_URL`",
            "`SEC_USER_AGENT`",
            "Optional `ALPHAVANTAGE_API_KEY`",
        ],
        size=13,
    )

    add_connector(slide, 3.9, 3.95, 4.18, 3.95, COLORS["teal"])
    add_connector(slide, 8.26, 3.95, 8.54, 3.95, COLORS["teal"])

    add_card(slide, 1.2, 6.0, 11.25, 0.58, COLORS["amber_bg"], COLORS["line"])
    add_text(
        slide,
        1.44,
        6.18,
        10.8,
        0.16,
        "Key simplification: the product can still run in public-data mode without Alpha Vantage, which reduces the risk of demo failure from free-tier rate limits.",
        size=12,
        color=COLORS["ink"],
    )
    add_slide_number(slide)


def build_close_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])
    add_title(slide, "Why this design works for the project", "A narrow scope, but a complete product loop.", dark=True)

    add_card(slide, 0.85, 1.8, 3.8, 3.9, COLORS["navy_soft"])
    add_text(slide, 1.12, 2.12, 3.0, 0.24, "What the project demonstrates", font_name=TITLE_FONT, size=20, color=COLORS["gold_soft"], bold=True)
    add_bullets(
        slide,
        1.12,
        2.62,
        3.05,
        [
            "LLM-based interaction, not static templates",
            "Tool and skill orchestration over public finance data",
            "Visible agent outputs beyond chat text",
            "Public URL that can be tested by the instructor",
        ],
        size=13,
        color=COLORS["white"],
    )

    add_card(slide, 4.95, 1.8, 3.8, 3.9, COLORS["navy_soft"])
    add_text(slide, 5.22, 2.12, 3.1, 0.24, "Operational choices", font_name=TITLE_FONT, size=20, color=COLORS["gold_soft"], bold=True)
    add_bullets(
        slide,
        5.22,
        2.62,
        3.05,
        [
            "Prompt-only UX to feel like an agent",
            "Public-data mode for deployment simplicity",
            "Fallbacks to avoid live demo collapse",
            "Research-only framing for safer positioning",
        ],
        size=13,
        color=COLORS["white"],
    )

    add_card(slide, 9.05, 1.8, 3.4, 3.9, COLORS["navy_soft"])
    add_text(slide, 9.32, 2.12, 2.8, 0.24, "Demo prompts", font_name=TITLE_FONT, size=20, color=COLORS["gold_soft"], bold=True)
    add_bullets(
        slide,
        9.32,
        2.62,
        2.7,
        [
            "Summarize Apple and show the recent price trend.",
            "Compare AMD and Nvidia on valuation and performance.",
            "Give me a market snapshot for SPY and highlight the risks.",
        ],
        size=12,
        color=COLORS["white"],
    )

    add_text(slide, 0.88, 6.2, 6.0, 0.3, "Public URL: pycallingagent.onrender.com", size=14, color=COLORS["white"], bold=True)
    add_text(slide, 0.88, 6.54, 6.2, 0.2, "Repository: github.com/vanzll/PyAgent", size=12, color=RGBColor(0xD4, 0xDE, 0xEC))
    add_footer(slide, "For research use only. Not investment advice.", dark=True)
    add_slide_number(slide, dark=True)


build_title_slide()
build_design_principles_slide()
build_architecture_slide()
build_execution_slide()
build_deployment_slide()
build_close_slide()

OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PATH)
print(f"wrote {OUT_PATH}")
