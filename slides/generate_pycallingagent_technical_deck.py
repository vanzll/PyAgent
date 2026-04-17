"""Generate the PyCallingAgent technical deck (~10 slides).

Run:
    python slides/generate_pycallingagent_technical_deck.py

Output: slides/pycallingagent_technical_deck.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


OUT_PATH = (
    Path(sys.argv[1])
    if len(sys.argv) > 1
    else Path(__file__).with_name("pycallingagent_technical_deck.pptx")
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
prs.core_properties.author = "Zhenglin Wan"
prs.core_properties.title = "PyCallingAgent - Technical Talk"
prs.core_properties.subject = "PyCallingAgent: runtime, skills, web product"

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"
MONO_FONT = "Consolas"

# --- Modern palette ---------------------------------------------------------
COLORS = {
    # primary surfaces
    "ink": RGBColor(0x10, 0x18, 0x28),
    "navy": RGBColor(0x0B, 0x1C, 0x32),
    "navy_soft": RGBColor(0x17, 0x2A, 0x45),
    "paper": RGBColor(0xFB, 0xF8, 0xF2),
    "sand": RGBColor(0xF3, 0xEE, 0xE2),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    # accents
    "amber": RGBColor(0xE0, 0x8A, 0x1F),
    "amber_soft": RGBColor(0xF6, 0xD9, 0xA0),
    "amber_bg": RGBColor(0xFB, 0xF1, 0xDD),
    "teal": RGBColor(0x0E, 0x6B, 0x66),
    "teal_soft": RGBColor(0xCC, 0xE6, 0xE2),
    "teal_bg": RGBColor(0xE3, 0xF2, 0xEE),
    "blue": RGBColor(0x2B, 0x4F, 0x9E),
    "blue_bg": RGBColor(0xE5, 0xEC, 0xF7),
    "rose_bg": RGBColor(0xFB, 0xEA, 0xE6),
    "muted": RGBColor(0x5B, 0x5F, 0x6B),
    "line": RGBColor(0xD7, 0xCF, 0xC2),
    "gray_bg": RGBColor(0xEE, 0xEB, 0xE3),
    # dark surface text
    "paper_on_dark": RGBColor(0xE9, 0xEC, 0xF3),
    "muted_on_dark": RGBColor(0xB6, 0xBD, 0xCB),
}

slide_no = 0


# --- Helpers ----------------------------------------------------------------
def add_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    x, y, w, h,
    text,
    *,
    font_name=BODY_FONT,
    size=14,
    color=None,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_card(slide, x, y, w, h, fill, line=None, line_width=Pt(1), radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_accent_bar(slide, x, y, color, *, length=0.5, weight=Pt(3)):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + length), Inches(y)
    )
    line.line.color.rgb = color
    line.line.width = weight
    return line


def add_vertical_bar(slide, x, y, h, color, *, weight=Pt(3)):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + h)
    )
    line.line.color.rgb = color
    line.line.width = weight
    return line


def add_pill(slide, x, y, w, text, fill, color, *, size=10, h=0.34):
    add_card(slide, x, y, w, h, fill, fill)
    add_text(
        slide, x, y, w, h, text,
        size=size, color=color, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def add_title(slide, eyebrow: str, title: str, subtitle: str | None = None, *, dark=False):
    eye_color = COLORS["amber"] if dark else COLORS["amber"]
    title_color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = COLORS["muted_on_dark"] if dark else COLORS["muted"]
    add_text(slide, 0.6, 0.45, 10.0, 0.22, eyebrow.upper(), size=10, color=eye_color, bold=True)
    add_text(
        slide, 0.6, 0.73, 12.0, 0.58, title,
        font_name=TITLE_FONT, size=28, color=title_color, bold=True,
    )
    if subtitle:
        add_text(slide, 0.6, 1.30, 12.0, 0.3, subtitle, size=13, color=sub_color, italic=True)
    # thin separator
    sep = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(0.6), Inches(1.66), Inches(12.7), Inches(1.66)
    )
    sep.line.color.rgb = COLORS["line"] if not dark else RGBColor(0x39, 0x46, 0x5E)
    sep.line.width = Pt(0.75)


def add_bullets(slide, x, y, w, items, *, size=13, color=None, line_spacing=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4.2))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(line_spacing)
    return box


def add_slide_number(slide, *, dark=False):
    global slide_no
    slide_no += 1
    add_text(
        slide, 12.6, 7.08, 0.4, 0.2,
        f"{slide_no:02d}", size=10,
        color=COLORS["muted_on_dark"] if dark else COLORS["muted"],
        align=PP_ALIGN.RIGHT, bold=True,
    )


def add_footer_brand(slide, *, dark=False):
    color = COLORS["muted_on_dark"] if dark else COLORS["muted"]
    add_text(
        slide, 0.6, 7.08, 10.0, 0.2,
        "PyCallingAgent  ·  Financial Research Agent",
        size=10, color=color, italic=True,
    )


def add_connector(slide, x1, y1, x2, y2, color, *, weight=Pt(1.75)):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = weight
    return line


def add_code_block(slide, x, y, w, h, lines, *, size=10, fill=None, color=None):
    add_card(slide, x, y, w, h, fill or COLORS["navy"], COLORS["navy"])
    box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.18), Inches(w - 0.4), Inches(h - 0.3))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = MONO_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color or COLORS["paper_on_dark"]
        p.space_after = Pt(0)
    return box


# --- Slide 1: Title ---------------------------------------------------------
def build_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])

    # left accent column
    add_card(slide, 0.0, 0.0, 0.28, 7.5, COLORS["amber"], COLORS["amber"], radius=False)

    add_pill(
        slide, 0.9, 0.7, 2.4,
        "TECHNICAL TALK", COLORS["amber"], COLORS["ink"],
    )

    add_text(
        slide, 0.9, 1.45, 11.5, 1.0,
        "PyCallingAgent",
        font_name=TITLE_FONT, size=62, color=COLORS["white"], bold=True,
    )
    add_text(
        slide, 0.9, 2.62, 11.5, 0.5,
        "A financial research agent built on a stateful Python runtime",
        size=20, color=COLORS["paper_on_dark"], italic=True,
    )

    # short description card
    add_card(slide, 0.9, 3.6, 5.7, 1.5, COLORS["navy_soft"], COLORS["navy_soft"])
    add_text(slide, 1.1, 3.78, 5.3, 0.26, "What this talk covers", size=11, color=COLORS["amber_soft"], bold=True)
    add_text(
        slide, 1.1, 4.1, 5.3, 0.92,
        "The runtime, skills, and web-layer adaptation that turn an LLM into an operator of live Python state, grounded on SEC and public market data.",
        size=13, color=COLORS["paper_on_dark"],
    )

    # stats stacked
    stats = [
        ("Prompt-first", "single box of input"),
        ("Object-native", "real DataFrames, not text"),
        ("Artifact outputs", "charts, tables, reports"),
        ("Public demo", "pycallingagent.onrender.com"),
    ]
    x = 6.95
    for idx, (label, sub) in enumerate(stats):
        col = idx % 2
        row = idx // 2
        left = x + col * 2.8
        top = 3.6 + row * 0.8
        add_card(slide, left, top, 2.65, 0.7, COLORS["navy_soft"], COLORS["navy_soft"])
        add_accent_bar(slide, left + 0.18, top + 0.15, COLORS["amber"], length=0.3)
        add_text(slide, left + 0.18, top + 0.24, 2.4, 0.22, label, size=12, color=COLORS["white"], bold=True)
        add_text(slide, left + 0.18, top + 0.44, 2.4, 0.2, sub, size=9, color=COLORS["muted_on_dark"])

    # footer strip
    add_card(slide, 0.9, 5.7, 11.5, 0.06, COLORS["amber"], COLORS["amber"], radius=False)
    add_text(
        slide, 0.9, 5.85, 11.5, 0.3,
        "Built on the CaveAgent runtime and skills  ·  adapted into a public FastAPI web product",
        size=12, color=COLORS["paper_on_dark"],
    )
    add_text(
        slide, 0.9, 6.2, 11.5, 0.25,
        "Zhenglin Wan  ·  NUS  ·  github.com/vanzll/PyAgent",
        size=11, color=COLORS["muted_on_dark"],
    )
    add_text(
        slide, 0.9, 6.5, 11.5, 0.25,
        "Research-only demonstration. Not investment advice.",
        size=10, color=COLORS["muted_on_dark"], italic=True,
    )
    add_slide_number(slide, dark=True)


# --- Slide 2: Problem & motivation ------------------------------------------
def build_motivation_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["paper"])
    add_title(
        slide,
        "Problem & motivation",
        "Classic tool calling hits a text bottleneck",
        "JSON-style tool calls force rich objects through a string channel on every turn.",
    )

    # Left: illustrative JSON snippet
    add_text(slide, 0.6, 1.95, 5.6, 0.26, "Conventional JSON function call", size=12, color=COLORS["blue"], bold=True)
    add_code_block(
        slide, 0.6, 2.28, 5.6, 2.6,
        [
            "// typical tool call",
            "{",
            '  "name": "get_prices",',
            '  "arguments": { "ticker": "AAPL" }',
            "}",
            "",
            "-> result serialized back as a string",
            "-> no object survives across calls",
        ],
        size=12,
    )

    # three pain points
    pains = [
        ("Serialization tax", "Every DataFrame, chart, or model is flattened into a string before the LLM can touch it again.", COLORS["blue"], COLORS["blue_bg"]),
        ("No persistent state", "Each call starts fresh. Long-lived objects and open connections cannot carry across turns.", COLORS["amber"], COLORS["amber_bg"]),
        ("Weak composition", "The model cannot chain in-memory objects with simple code; it re-describes data instead of computing on it.", COLORS["teal"], COLORS["teal_bg"]),
    ]
    y = 2.0
    for title, body, accent, bg in pains:
        add_card(slide, 6.6, y, 6.2, 0.94, bg, COLORS["line"])
        add_vertical_bar(slide, 6.72, y + 0.12, 0.7, accent, weight=Pt(3))
        add_text(slide, 6.95, y + 0.14, 5.6, 0.26, title, size=13, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, 6.95, y + 0.42, 5.6, 0.5, body, size=11, color=COLORS["muted"])
        y += 1.04

    # bottom takeaway
    add_card(slide, 0.6, 5.55, 12.2, 1.2, COLORS["ink"], COLORS["ink"])
    add_accent_bar(slide, 0.85, 5.78, COLORS["amber"], length=0.45, weight=Pt(3))
    add_text(
        slide, 0.85, 5.9, 11.8, 0.3,
        "Takeaway",
        size=12, color=COLORS["amber_soft"], bold=True,
    )
    add_text(
        slide, 0.85, 6.18, 11.8, 0.5,
        "For real research workloads we want the LLM to drive a live Python process, not to replay text representations of it.",
        size=14, color=COLORS["paper_on_dark"],
    )
    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 3: Design philosophy / paradigm evolution ------------------------
def build_paradigm_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["sand"])
    add_title(
        slide,
        "Design philosophy",
        "From text emitter to process operator",
        "We position PyCallingAgent at the next step of the tool-calling timeline.",
    )

    # four-step paradigm timeline
    stages = [
        ("Gen 0", "Text only", "Plain chat. No tools.", COLORS["muted"], COLORS["gray_bg"]),
        ("Gen 1", "JSON FC", "Stateless function calls, string I/O.", COLORS["blue"], COLORS["blue_bg"]),
        ("Gen 2", "Code-as-action", "Scratch interpreter per call, still sandboxed text.", COLORS["amber"], COLORS["amber_bg"]),
        ("Gen 3", "Stateful runtime", "Persistent kernel; objects live across turns.", COLORS["teal"], COLORS["teal_bg"]),
    ]

    x0 = 0.6
    w = 3.06
    gap = 0.12
    y = 2.05
    h = 2.05
    for idx, (step, title, body, accent, bg) in enumerate(stages):
        x = x0 + idx * (w + gap)
        add_card(slide, x, y, w, h, bg, COLORS["line"])
        add_vertical_bar(slide, x + 0.14, y + 0.2, 0.7, accent, weight=Pt(3))
        add_text(slide, x + 0.4, y + 0.22, w - 0.5, 0.22, step, size=10, color=accent, bold=True)
        add_text(slide, x + 0.4, y + 0.5, w - 0.5, 0.3, title, size=17, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, x + 0.4, y + 0.98, w - 0.5, 1.0, body, size=11, color=COLORS["muted"])
        if idx < len(stages) - 1:
            add_connector(slide, x + w + 0.005, y + h / 2, x + w + gap - 0.005, y + h / 2, COLORS["amber"], weight=Pt(2))

    # pillar cards at bottom
    add_text(
        slide, 0.6, 4.5, 12.2, 0.3,
        "Three pillars that make Gen 3 practical",
        size=13, color=COLORS["ink"], bold=True, font_name=TITLE_FONT,
    )

    pillars = [
        ("Twin-track loop", "A reasoning track (LLM messages) paired with an execution track (persistent Python)."),
        ("Deterministic state access", "The runtime is inspectable: we can check variables, verify outputs, score episodes."),
        ("Runtime-aware skills", "Skills inject live callables instead of only text instructions."),
    ]
    x = 0.6
    for idx, (title, body) in enumerate(pillars):
        cx = x + idx * 4.15
        add_card(slide, cx, 4.88, 3.95, 1.72, COLORS["white"], COLORS["line"])
        add_accent_bar(slide, cx + 0.28, 5.06, COLORS["amber"], length=0.4)
        add_text(slide, cx + 0.28, 5.18, 3.5, 0.3, title, size=14, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, cx + 0.28, 5.58, 3.5, 1.0, body, size=11, color=COLORS["muted"])

    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 4: Core architecture (dual-stream) -------------------------------
def build_architecture_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["paper"])
    add_title(
        slide,
        "Core architecture",
        "Paired semantic and runtime tracks",
        "The agent loop binds an LLM message history to a persistent Python environment.",
    )

    # Top - semantic stream
    y_sem = 1.95
    add_card(slide, 0.6, y_sem, 12.2, 1.55, COLORS["blue_bg"], COLORS["line"])
    add_vertical_bar(slide, 0.72, y_sem + 0.2, 1.15, COLORS["blue"], weight=Pt(3))
    add_text(slide, 0.95, y_sem + 0.18, 4.5, 0.24, "Semantic track", size=12, color=COLORS["blue"], bold=True)
    add_text(slide, 0.95, y_sem + 0.46, 4.5, 0.28, "Reasoning & intent", size=16, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
    add_text(
        slide, 0.95, y_sem + 0.82, 4.5, 0.6,
        "Rolling LLM history. Each turn updates h_t = LLM(x_t, h_{t-1}).",
        size=11, color=COLORS["muted"],
    )

    # four mini LLM bubbles
    bx = 5.9
    for i in range(4):
        cx = bx + i * 1.65
        add_card(slide, cx, y_sem + 0.42, 1.3, 0.68, COLORS["white"], COLORS["blue"])
        add_text(slide, cx, y_sem + 0.5, 1.3, 0.24, f"h_{{t+{i}}}", size=12, color=COLORS["blue"], bold=True, align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
        add_text(slide, cx, y_sem + 0.74, 1.3, 0.22, "LLM step", size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if i < 3:
            add_connector(slide, cx + 1.3, y_sem + 0.76, cx + 1.65, y_sem + 0.76, COLORS["blue"])
    add_text(slide, 5.9, y_sem + 1.15, 6.6, 0.22, "text output -> user answer", size=10, color=COLORS["muted"], italic=True, align=PP_ALIGN.CENTER)

    # middle flow arrows
    for i in range(4):
        cx = bx + i * 1.65 + 0.65
        add_connector(slide, cx, y_sem + 1.42, cx, y_sem + 1.95, COLORS["amber"], weight=Pt(1.5))

    # Bottom - runtime stream
    y_rt = 3.95
    add_card(slide, 0.6, y_rt, 12.2, 1.85, COLORS["teal_bg"], COLORS["line"])
    add_vertical_bar(slide, 0.72, y_rt + 0.2, 1.45, COLORS["teal"], weight=Pt(3))
    add_text(slide, 0.95, y_rt + 0.18, 4.5, 0.24, "Runtime track", size=12, color=COLORS["teal"], bold=True)
    add_text(slide, 0.95, y_rt + 0.46, 4.5, 0.28, "Persistent state & data", size=16, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
    add_text(
        slide, 0.95, y_rt + 0.82, 4.5, 0.8,
        "Long-lived IPython kernel. State update rule: s_t = exec(c_t, s_{t-1}).",
        size=11, color=COLORS["muted"],
    )

    # state slots
    for i in range(4):
        cx = bx + i * 1.65
        add_card(slide, cx, y_rt + 0.46, 1.3, 0.75, COLORS["white"], COLORS["teal"])
        add_text(slide, cx, y_rt + 0.54, 1.3, 0.24, f"s_{{t+{i}}}", size=12, color=COLORS["teal"], bold=True, align=PP_ALIGN.CENTER, font_name=TITLE_FONT)
        add_text(slide, cx, y_rt + 0.82, 1.3, 0.22, "DataFrames", size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        add_text(slide, cx, y_rt + 1.0, 1.3, 0.22, "objects · funcs", size=9, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        if i < 3:
            add_connector(slide, cx + 1.3, y_rt + 0.82, cx + 1.65, y_rt + 0.82, COLORS["teal"])

    add_text(
        slide, 0.95, y_rt + 1.42, 11.6, 0.3,
        "Observation shaping: only compact summaries are folded back into the prompt; full objects stay in the kernel.",
        size=11, color=COLORS["muted"], italic=True,
    )

    # bottom caption strip
    add_card(slide, 0.6, 6.05, 12.2, 0.7, COLORS["ink"], COLORS["ink"])
    add_accent_bar(slide, 0.85, 6.25, COLORS["amber"], length=0.45)
    add_text(
        slide, 0.85, 6.33, 11.8, 0.32,
        "The LLM context stays lean. The runtime is the real workspace.",
        size=13, color=COLORS["paper_on_dark"], bold=True,
    )
    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 5: Runtime API + skills ------------------------------------------
def build_runtime_skills_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["sand"])
    add_title(
        slide,
        "Runtime & skills",
        "Injection API plus lazy skill loading",
        "Objects live in the kernel; the LLM only sees compact descriptors.",
    )

    # left: code
    add_text(slide, 0.6, 1.95, 6.4, 0.26, "Runtime injection (Python)", size=12, color=COLORS["teal"], bold=True)
    add_code_block(
        slide, 0.6, 2.28, 6.4, 3.2,
        [
            "from cave_agent import CaveAgent, Variable, Function",
            "from cave_agent.runtime import IPythonRuntime",
            "",
            "runtime = IPythonRuntime()",
            "runtime.inject_variable(Variable(",
            "    name='prices',",
            "    value=price_df,       # real pandas DataFrame",
            "    description='daily OHLCV for AAPL',",
            "))",
            "runtime.inject_function(Function(fetch_fundamentals))",
            "",
            "agent = CaveAgent(model=model, runtime=runtime)",
            "await agent.run('Compare AAPL and NVDA on valuation.')",
        ],
        size=11,
    )

    # right: bullets & skill mechanism
    add_text(slide, 7.3, 1.95, 5.5, 0.26, "Runtime primitives", size=12, color=COLORS["teal"], bold=True)

    prims = [
        ("Variable", "any Python object, retrievable by name."),
        ("Function", "a callable surfaced to the LLM with a docstring."),
        ("Type", "schema auto-extracted from Pydantic, dataclass, or Enum."),
    ]
    y = 2.28
    for name, body in prims:
        add_card(slide, 7.3, y, 5.5, 0.66, COLORS["white"], COLORS["line"])
        add_text(slide, 7.5, y + 0.12, 1.6, 0.22, name, size=12, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, 8.95, y + 0.14, 3.75, 0.45, body, size=10, color=COLORS["muted"])
        y += 0.76

    # skills panel
    add_card(slide, 7.3, 4.6, 5.5, 1.9, COLORS["teal_bg"], COLORS["line"])
    add_accent_bar(slide, 7.48, 4.78, COLORS["teal"], length=0.45)
    add_text(slide, 7.48, 4.88, 5.2, 0.26, "Skill mechanism", size=13, color=COLORS["teal"], bold=True, font_name=TITLE_FONT)
    add_bullets(
        slide, 7.48, 5.2, 5.1,
        [
            "Lazy loading: only skill metadata (~100 tokens) at startup.",
            "activate_skill(name) injects full instructions and callables on demand.",
            "Each skill can export Function / Variable / Type via injection.py.",
        ],
        size=10, color=COLORS["ink"], line_spacing=4,
    )

    add_card(slide, 0.6, 5.7, 6.4, 1.0, COLORS["ink"], COLORS["ink"])
    add_accent_bar(slide, 0.85, 5.88, COLORS["amber"], length=0.4)
    add_text(
        slide, 0.85, 5.98, 6.1, 0.22,
        "Context budget saved",
        size=11, color=COLORS["amber_soft"], bold=True,
    )
    add_text(
        slide, 0.85, 6.24, 6.1, 0.7,
        "Variables, functions, types, and skills are advertised as short descriptors in the prompt.",
        size=10, color=COLORS["paper_on_dark"],
    )
    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 6: Financial research workflow -----------------------------------
def build_workflow_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["paper"])
    add_title(
        slide,
        "Financial research workflow",
        "What happens when the user presses Send",
        "One prompt triggers a full research pass with grounded outputs.",
    )

    steps = [
        ("01", "Parse prompt", "Route finance / generic; extract tickers or market proxies.", COLORS["blue"], COLORS["blue_bg"]),
        ("02", "Fetch grounded data", "SEC EDGAR filings, optional Alpha Vantage, public fallback.", COLORS["teal"], COLORS["teal_bg"]),
        ("03", "Inject into runtime", "Price frames, fundamentals, and helper functions become live variables.", COLORS["amber"], COLORS["amber_bg"]),
        ("04", "Run agent loop", "The LLM calls skills over real objects; no text-only reasoning about the numbers.", COLORS["blue"], COLORS["rose_bg"]),
        ("05", "Materialize outputs", "Save charts, tables, a markdown report, and snapshot cards to the session.", COLORS["teal"], COLORS["gray_bg"]),
    ]

    w = 2.34
    gap = 0.14
    x0 = 0.6
    y = 2.05
    h = 3.5
    for idx, (step, title, body, accent, bg) in enumerate(steps):
        x = x0 + idx * (w + gap)
        add_card(slide, x, y, w, h, bg, COLORS["line"])
        add_vertical_bar(slide, x + 0.18, y + 0.25, 0.75, accent, weight=Pt(3))
        add_text(slide, x + 0.45, y + 0.28, w - 0.6, 0.3, step, size=16, color=accent, bold=True, font_name=TITLE_FONT)
        add_text(slide, x + 0.45, y + 0.78, w - 0.6, 0.3, title, size=14, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, x + 0.45, y + 1.22, w - 0.6, 2.0, body, size=11, color=COLORS["muted"])

    # connector arrows between steps
    for i in range(4):
        cx = x0 + i * (w + gap) + w
        add_connector(slide, cx - 0.02, y + h / 2, cx + gap + 0.02, y + h / 2, COLORS["amber"], weight=Pt(2))

    # two-band caption strip
    add_card(slide, 0.6, 5.85, 12.2, 0.85, COLORS["ink"], COLORS["ink"])
    add_accent_bar(slide, 0.85, 6.1, COLORS["amber"], length=0.45)
    add_text(
        slide, 0.85, 6.18, 11.8, 0.28,
        "One turn = both language answer and visible work product",
        size=13, color=COLORS["paper_on_dark"], bold=True,
    )
    add_text(
        slide, 0.85, 6.44, 11.8, 0.22,
        "Non-financial prompts fall through to generic LLM QA without touching the finance pipeline.",
        size=10, color=COLORS["muted_on_dark"], italic=True,
    )
    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 7: Application-layer adaptation ----------------------------------
def build_app_layer_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["sand"])
    add_title(
        slide,
        "Application-layer adaptation",
        "Wrapping the agent as a FastAPI product",
        "The agent core stays the same; the web layer adds session/run orchestration and streaming.",
    )

    # Architecture row
    y = 2.0
    nodes = [
        ("Browser UI", "Jinja2 + vanilla JS\nprompt box, chat, artifact panel", COLORS["blue_bg"], COLORS["blue"]),
        ("FastAPI service", "Sessions, runs, SSE/polling, file uploads\ndownload endpoints per artifact", COLORS["amber_bg"], COLORS["amber"]),
        ("Agent runner", "FinancialResearchRunner\ninfer tickers -> fetch data -> run agent -> save outputs", COLORS["teal_bg"], COLORS["teal"]),
        ("Data layer", "SEC EDGAR + optional Alpha Vantage\npublic-data mode as safe fallback", COLORS["rose_bg"], COLORS["blue"]),
    ]
    w = 2.94
    gap = 0.13
    h = 1.78
    for idx, (title, body, bg, accent) in enumerate(nodes):
        x = 0.6 + idx * (w + gap)
        add_card(slide, x, y, w, h, bg, COLORS["line"])
        add_vertical_bar(slide, x + 0.16, y + 0.22, 0.6, accent, weight=Pt(3))
        add_text(slide, x + 0.4, y + 0.22, w - 0.6, 0.28, title, size=14, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, x + 0.4, y + 0.7, w - 0.6, 1.0, body, size=11, color=COLORS["muted"])
        if idx < len(nodes) - 1:
            add_connector(slide, x + w - 0.02, y + h / 2, x + w + gap + 0.02, y + h / 2, COLORS["amber"], weight=Pt(2))

    # Three detail cards
    details = [
        (
            "Session + run orchestration",
            [
                "Each session owns a message history and a working folder.",
                "Every user prompt becomes a run with its own event log and artifact set.",
                "Runs are cancellable; the kernel is reusable across turns.",
            ],
            COLORS["amber"],
        ),
        (
            "Live updates to the page",
            [
                "SSE-style event stream plus polling fallback for cold networks.",
                "Sliding window of recent events drives the Live Status panel.",
                "Friendly error surface instead of a generic failure banner.",
            ],
            COLORS["teal"],
        ),
        (
            "Artifact rendering",
            [
                "Tables and charts preview inline; full files are downloadable.",
                "Markdown report rendered in-page after each successful run.",
                "Non-financial prompts bypass the pipeline and chat normally.",
            ],
            COLORS["blue"],
        ),
    ]
    y = 4.0
    for idx, (title, bullets, accent) in enumerate(details):
        x = 0.6 + idx * 4.14
        add_card(slide, x, y, 3.94, 2.45, COLORS["white"], COLORS["line"])
        add_accent_bar(slide, x + 0.25, y + 0.2, accent, length=0.45, weight=Pt(3))
        add_text(slide, x + 0.25, y + 0.32, 3.5, 0.28, title, size=13, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_bullets(slide, x + 0.25, y + 0.68, 3.5, bullets, size=10, color=COLORS["muted"], line_spacing=4)

    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 8: Deployment ----------------------------------------------------
def build_deployment_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["paper"])
    add_title(
        slide,
        "Deployment",
        "From GitHub to a public URL in one path",
        "The deployment surface is small on purpose so the demo does not break.",
    )

    # Three big column cards
    cols = [
        (
            "GitHub",
            "Source of truth",
            [
                "Push to `main` on vanzll/PyAgent.",
                "Render watches the branch and redeploys automatically.",
            ],
            COLORS["blue_bg"], COLORS["blue"],
        ),
        (
            "Render",
            "Web service",
            [
                "Build: `pip install \".[openai,webapp]\"`",
                "Start: `python -m uvicorn cave_agent.webapp.app:create_app --factory`",
                "Public URL comes from Render; no custom DNS needed.",
            ],
            COLORS["amber_bg"], COLORS["amber"],
        ),
        (
            "Runtime config",
            "Environment variables",
            [
                "`LLM_MODEL_ID`, `LLM_API_KEY`, `LLM_BASE_URL`",
                "`SEC_USER_AGENT` for SEC EDGAR traffic.",
                "`ALPHAVANTAGE_API_KEY` is optional; public-data mode covers the fallback.",
            ],
            COLORS["teal_bg"], COLORS["teal"],
        ),
    ]
    y = 2.0
    for idx, (title, sub, bullets, bg, accent) in enumerate(cols):
        x = 0.6 + idx * 4.15
        add_card(slide, x, y, 3.95, 3.9, bg, COLORS["line"])
        add_accent_bar(slide, x + 0.3, y + 0.25, accent, length=0.45, weight=Pt(3))
        add_text(slide, x + 0.3, y + 0.38, 3.4, 0.3, title, size=18, color=COLORS["ink"], bold=True, font_name=TITLE_FONT)
        add_text(slide, x + 0.3, y + 0.76, 3.4, 0.24, sub, size=11, color=accent, bold=True)
        add_bullets(slide, x + 0.3, y + 1.12, 3.5, bullets, size=11, color=COLORS["muted"], line_spacing=6)
        if idx < len(cols) - 1:
            add_connector(slide, x + 3.95, y + 1.95, x + 4.15, y + 1.95, COLORS["amber"], weight=Pt(2))

    # Bottom strip: public URL + notes
    add_card(slide, 0.6, 6.1, 12.2, 0.95, COLORS["ink"], COLORS["ink"])
    add_accent_bar(slide, 0.85, 6.32, COLORS["amber"], length=0.45, weight=Pt(3))
    add_text(
        slide, 0.85, 6.42, 11.8, 0.28,
        "Live demo: pycallingagent.onrender.com",
        size=14, color=COLORS["paper_on_dark"], bold=True,
    )
    add_text(
        slide, 0.85, 6.72, 11.8, 0.22,
        "Public-data fallback keeps the deployment working even if third-party rate limits kick in.",
        size=10, color=COLORS["muted_on_dark"], italic=True,
    )
    add_footer_brand(slide)
    add_slide_number(slide)


# --- Slide 9: Demo placeholder 1 --------------------------------------------
def build_demo_slide(idx: int, title: str, subtitle: str, hint: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLORS["navy"])
    add_card(slide, 0.0, 0.0, 0.28, 7.5, COLORS["amber"], COLORS["amber"], radius=False)

    add_pill(slide, 0.9, 0.7, 1.7, f"DEMO {idx}", COLORS["amber"], COLORS["ink"])
    add_text(
        slide, 0.9, 1.35, 11.8, 0.8,
        title, font_name=TITLE_FONT, size=40, color=COLORS["white"], bold=True,
    )
    add_text(
        slide, 0.9, 2.45, 11.8, 0.4,
        subtitle, size=16, color=COLORS["paper_on_dark"], italic=True,
    )

    # large placeholder frame
    add_card(slide, 0.9, 3.25, 11.6, 3.3, COLORS["navy_soft"], COLORS["amber"], line_width=Pt(1.5))
    add_text(
        slide, 0.9, 4.15, 11.6, 0.5,
        "[ place demo screenshot or live screen here ]",
        size=18, color=COLORS["muted_on_dark"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, 0.9, 4.7, 11.6, 0.4,
        hint, size=12, color=COLORS["amber_soft"],
        align=PP_ALIGN.CENTER, italic=True,
    )
    add_text(
        slide, 0.9, 6.75, 11.8, 0.25,
        "pycallingagent.onrender.com  ·  github.com/vanzll/PyAgent",
        size=10, color=COLORS["muted_on_dark"],
    )
    add_slide_number(slide, dark=True)


# --- Build all --------------------------------------------------------------
build_title_slide()
build_motivation_slide()
build_paradigm_slide()
build_architecture_slide()
build_runtime_skills_slide()
build_workflow_slide()
build_app_layer_slide()
build_deployment_slide()
build_demo_slide(
    1,
    "Demo",
    "Prompt -> grounded research output",
    "Suggested prompt:  Summarize Apple and show the recent price trend.",
)
build_demo_slide(
    2,
    "Demo",
    "Comparison and artifacts",
    "Suggested prompt:  Compare AMD and Nvidia on valuation and performance.",
)

OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PATH)
print(f"wrote {OUT_PATH}")
